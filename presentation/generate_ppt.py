from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


OUTPUT = Path(__file__).resolve().parent / "Spring-Phone-Presentation.pptx"

BG = RGBColor(255, 249, 245)
TEXT = RGBColor(63, 54, 82)
SUBTEXT = RGBColor(107, 100, 128)
WHITE = RGBColor(255, 255, 255)
MINT = RGBColor(190, 234, 215)
BLUE = RGBColor(191, 223, 255)
YELLOW = RGBColor(255, 241, 184)
PEACH = RGBColor(255, 214, 201)
LAVENDER = RGBColor(233, 228, 245)
ROSE = RGBColor(255, 241, 247)
GRAY = RGBColor(230, 226, 234)
GREEN = RGBColor(110, 176, 126)


SLIDES = [
    {
        "kind": "cover",
        "title": "Spring Phone",
        "subtitle": "A Personalized App Store Backend Built with Spring Boot",
        "tagline": "Browse it. Love it. Buy it. Keep it.",
    },
    {
        "kind": "hero",
        "tag": "Product Vision",
        "title": "What Problem Are We Solving?",
        "subtitle": "A tiny app store experience that feels personal, friendly, and organized.",
        "points": [
            "Explore apps easily",
            "Purchase apps smoothly",
            "See what you already own",
        ],
        "note": "After a user buys an app, that app disappears from that user's available list.",
    },
    {
        "kind": "journey",
        "tag": "User Journey",
        "title": "Meet Alice",
        "subtitle": "One cute user story guides the whole project.",
        "steps": [
            ("Browse", "Alice opens the store and looks around."),
            ("Filter", "She narrows apps by name or category."),
            ("Purchase", "She buys one app she really likes."),
            ("My Apps", "She now sees it in her purchased list."),
        ],
    },
    {
        "kind": "journey",
        "tag": "Business Logic",
        "title": "What Changes After Purchase?",
        "subtitle": "The store view becomes personal for Alice.",
        "steps": [
            ("Before", "TravelMap appears in available apps."),
            ("Buy", "Alice purchases TravelMap."),
            ("After", "TravelMap moves into her owned apps."),
            ("Result", "It no longer appears in her available list."),
        ],
    },
    {
        "kind": "architecture",
        "tag": "Architecture",
        "title": "Three-Layer Architecture",
        "subtitle": "A classic structure, organized around the user journey.",
        "layers": [
            ("Controller Layer", "Handles HTTP requests and responses"),
            ("Service Layer", "Contains business rules and core logic"),
            ("Repository Layer", "Reads and writes data"),
            ("MySQL Database", "Stores apps, users, and purchases"),
        ],
    },
    {
        "kind": "split",
        "tag": "Design Choice",
        "title": "Why Three Layers?",
        "subtitle": "It keeps the system simple for today and flexible for tomorrow.",
        "left_title": "Benefits",
        "left_points": [
            "Clear responsibilities",
            "Easier maintenance",
            "Cleaner code",
        ],
        "right_title": "What It Means Here",
        "right_points": [
            "Controllers expose APIs",
            "Services enforce rules",
            "Repositories access MySQL",
        ],
    },
    {
        "kind": "domain",
        "tag": "Domain Model",
        "title": "Our Core Data Model",
        "subtitle": "Three simple entities build the whole app store story.",
        "cards": [
            ("App", "id, name, category, price, description"),
            ("User", "id, username, email"),
            ("Purchase", "id, user, app, purchaseTime"),
        ],
        "footer": "Purchase is the key entity because it records ownership.",
    },
    {
        "kind": "formula",
        "tag": "Personalization",
        "title": "What Makes the Store Personal?",
        "subtitle": "Available apps are calculated for each user.",
        "formula": "All Apps  -  Purchased Apps  =  Available Apps",
        "points": [
            "An app does not disappear from the whole store",
            "It only disappears for the user who already bought it",
            "Each user gets a personalized store view",
        ],
    },
    {
        "kind": "flow",
        "tag": "Request Flow",
        "title": "Example Flow: Alice Buys TravelMap",
        "subtitle": "One request moving through all three layers.",
        "steps": [
            "POST /api/users/1/purchases",
            "Controller receives the request",
            "Service checks the user and app",
            "Service validates duplicate purchase rules",
            "Repository saves the purchase",
            "TravelMap appears in My Apps and disappears from Available Apps",
        ],
    },
    {
        "kind": "api",
        "tag": "REST API",
        "title": "Browse and Discover APIs",
        "subtitle": "The first part of the journey is exploration.",
        "apis": [
            ("GET /api/apps", "Browse all apps"),
            ("GET /api/apps?name=...&category=...", "Search and filter apps"),
            ("GET /api/users/{userId}/apps/available", "Show apps available to one user"),
        ],
    },
    {
        "kind": "api",
        "tag": "REST API",
        "title": "Purchase and Ownership APIs",
        "subtitle": "The second part of the journey is ownership.",
        "apis": [
            ("POST /api/users/{userId}/purchases", "Purchase an app"),
            ("GET /api/users/{userId}/purchases", "List purchased apps"),
        ],
    },
    {
        "kind": "split",
        "tag": "Service Design",
        "title": "Where the Business Rules Live",
        "subtitle": "Our services are the brain of the system.",
        "left_title": "AppService",
        "left_points": [
            "Browse all apps",
            "Search by name",
            "Filter by category",
        ],
        "right_title": "PurchaseService",
        "right_points": [
            "Validate user and app",
            "Prevent duplicate purchase",
            "Return purchased apps",
        ],
    },
    {
        "kind": "split",
        "tag": "Persistence",
        "title": "How Data Lives Between Runs",
        "subtitle": "We connect the backend to MySQL so data stays safe after restart.",
        "left_title": "Stored Tables",
        "left_points": [
            "mobile_app",
            "store_user",
            "purchase",
        ],
        "right_title": "Why It Matters",
        "right_points": [
            "Real persistence",
            "Realistic demo data",
            "Reliable hackathon showcase",
        ],
    },
    {
        "kind": "testing",
        "tag": "Testing",
        "title": "How We Tested the Project",
        "subtitle": "We covered both happy paths and edge cases.",
        "good": [
            "Browse apps",
            "Filter by name and category",
            "Successful purchase",
            "Purchased list retrieval",
        ],
        "bad": [
            "Duplicate purchase",
            "Missing user",
            "Missing app",
            "Empty result cases",
        ],
    },
    {
        "kind": "future",
        "tag": "Future Ideas",
        "title": "If We Had More Time...",
        "subtitle": "A small project with room to grow.",
        "product": [
            "Ratings and reviews",
            "Recommendations",
            "Authentication",
            "Admin dashboard",
        ],
        "tech": [
            "Local cache for repeated purchase checks",
            "Pagination and sorting",
            "Global exception handling",
            "Integration tests",
        ],
    },
    {
        "kind": "closing",
        "tag": "Closing",
        "title": "Thank You",
        "subtitle": "From discovery to ownership, every step is simple, personal, and structured.",
        "points": [
            "We designed Spring Phone around a user journey",
            "We implemented it with a clean three-layer Spring architecture",
            "We built a backend that is small, clear, and ready to grow",
        ],
    },
]


def set_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def add_round_box(slide, left, top, width, height, fill_rgb, line_rgb=None):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    shape.line.color.rgb = fill_rgb if line_rgb is None else line_rgb
    return shape


def add_text(shape, text, size=18, bold=False, color=TEXT, align=PP_ALIGN.LEFT, font_name="Aptos"):
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tf


def add_bullet_list(shape, items, size=18, color=TEXT):
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {item}"
        p.font.name = "Aptos"
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(7)
    return tf


def add_decorations(slide):
    add_round_box(slide, Inches(8.15), Inches(0.3), Inches(1.25), Inches(0.18), MINT)
    add_round_box(slide, Inches(8.55), Inches(6.88), Inches(0.68), Inches(0.16), PEACH)
    add_round_box(slide, Inches(0.45), Inches(6.78), Inches(0.92), Inches(0.16), LAVENDER)


def add_header(slide, title, subtitle, tag=None):
    if tag:
        tag_box = add_round_box(slide, Inches(0.55), Inches(0.32), Inches(1.7), Inches(0.38), YELLOW)
        add_text(tag_box, tag, size=13, bold=True, align=PP_ALIGN.CENTER)
        tag_box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    title_box = add_round_box(slide, Inches(0.55), Inches(0.82), Inches(8.1), Inches(0.82), ROSE)
    add_text(title_box, title, size=24, bold=True, font_name="Aptos Display")

    subtitle_box = add_round_box(slide, Inches(0.85), Inches(1.62), Inches(7.5), Inches(0.54), RGBColor(244, 249, 255))
    tf = add_text(subtitle_box, subtitle, size=15, color=SUBTEXT)
    tf.paragraphs[0].runs[0].font.italic = True


def add_phone(slide, left, top, width=1.55, height=3.0):
    phone = add_round_box(slide, Inches(left), Inches(top), Inches(width), Inches(height), WHITE, LAVENDER)
    phone.line.width = Pt(1.2)
    screen = add_round_box(
        slide, Inches(left + 0.18), Inches(top + 0.28), Inches(width - 0.36), Inches(height - 0.52), RGBColor(248, 252, 255), BLUE
    )
    return phone, screen


def add_avatar(slide, left, top, scale=1.0):
    head = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(left), Inches(top), Inches(0.9 * scale), Inches(0.9 * scale))
    head.fill.solid()
    head.fill.fore_color.rgb = PEACH
    head.line.color.rgb = PEACH

    hair = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ARC, Inches(left - 0.02), Inches(top - 0.08), Inches(0.95 * scale), Inches(0.55 * scale))
    hair.fill.solid()
    hair.fill.fore_color.rgb = TEXT
    hair.line.color.rgb = TEXT

    body = add_round_box(slide, Inches(left - 0.06), Inches(top + 0.82 * scale), Inches(1.0 * scale), Inches(1.15 * scale), MINT, MINT)
    body.line.width = Pt(1.0)

    eye1 = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(left + 0.22), Inches(top + 0.34), Inches(0.06), Inches(0.06))
    eye1.fill.solid()
    eye1.fill.fore_color.rgb = TEXT
    eye1.line.color.rgb = TEXT
    eye2 = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(left + 0.54), Inches(top + 0.34), Inches(0.06), Inches(0.06))
    eye2.fill.solid()
    eye2.fill.fore_color.rgb = TEXT
    eye2.line.color.rgb = TEXT

    smile = slide.shapes.add_textbox(Inches(left + 0.18), Inches(top + 0.46), Inches(0.5), Inches(0.2))
    tf = smile.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "◡"
    run.font.size = Pt(16)
    run.font.color.rgb = TEXT

    label = add_round_box(slide, Inches(left - 0.12), Inches(top + 2.02 * scale), Inches(1.15 * scale), Inches(0.34), YELLOW)
    add_text(label, "Alice", size=12, bold=True, align=PP_ALIGN.CENTER)


def add_card(slide, left, top, width, height, title, body, color):
    box = add_round_box(slide, Inches(left), Inches(top), Inches(width), Inches(height), color)
    box.line.color.rgb = color
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    p1 = tf.paragraphs[0]
    r1 = p1.add_run()
    r1.text = title
    r1.font.name = "Aptos Display"
    r1.font.size = Pt(17)
    r1.font.bold = True
    r1.font.color.rgb = TEXT
    p2 = tf.add_paragraph()
    p2.text = body
    p2.font.name = "Aptos"
    p2.font.size = Pt(12)
    p2.font.color.rgb = TEXT
    return box


def render_cover(slide, data):
    add_decorations(slide)
    title_box = add_round_box(slide, Inches(0.65), Inches(1.1), Inches(5.2), Inches(1.0), ROSE)
    add_text(title_box, data["title"], size=28, bold=True, font_name="Aptos Display")

    subtitle_box = add_round_box(slide, Inches(0.95), Inches(2.12), Inches(4.8), Inches(0.7), RGBColor(244, 249, 255))
    tf = add_text(subtitle_box, data["subtitle"], size=16, color=SUBTEXT)
    tf.paragraphs[0].runs[0].font.italic = True

    tag = add_round_box(slide, Inches(0.95), Inches(2.95), Inches(3.2), Inches(0.58), YELLOW)
    add_text(tag, data["tagline"], size=16, bold=True, align=PP_ALIGN.CENTER)
    tag.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    for i, (x, y, c, label) in enumerate([
        (1.0, 4.1, MINT, "Spring Boot"),
        (2.45, 4.1, BLUE, "MySQL"),
        (3.9, 4.1, PEACH, "REST API"),
    ]):
        card = add_round_box(slide, Inches(x), Inches(y), Inches(1.15), Inches(1.0), c)
        add_text(card, label, size=12, bold=True, align=PP_ALIGN.CENTER)
        card.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    add_phone(slide, 6.7, 1.65, 2.0, 3.95)
    for x, y, c in [(6.98, 2.1, MINT), (7.55, 2.1, PEACH), (6.98, 2.78, YELLOW), (7.55, 2.78, LAVENDER)]:
        icon = add_round_box(slide, Inches(x), Inches(y), Inches(0.42), Inches(0.42), c)
        icon.line.color.rgb = c
    add_avatar(slide, 8.0, 4.45, 0.9)


def render_hero(slide, data):
    add_decorations(slide)
    add_header(slide, data["title"], data["subtitle"], data.get("tag"))
    note = add_round_box(slide, Inches(0.85), Inches(5.75), Inches(7.6), Inches(0.75), YELLOW)
    add_text(note, data["note"], size=15, bold=True, align=PP_ALIGN.CENTER)
    note.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    for i, point in enumerate(data["points"]):
        color = [MINT, BLUE, PEACH][i % 3]
        box = add_round_box(slide, Inches(0.95 + i * 2.45), Inches(2.6), Inches(2.05), Inches(2.45), color)
        add_text(box, point, size=18, bold=True, align=PP_ALIGN.CENTER)
        box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    add_phone(slide, 8.45, 2.4, 1.1, 2.25)


def render_journey(slide, data):
    add_decorations(slide)
    add_header(slide, data["title"], data["subtitle"], data.get("tag"))
    add_avatar(slide, 0.8, 2.45, 1.0)
    colors = [MINT, BLUE, YELLOW, PEACH]
    x = 2.05
    for i, (title, body) in enumerate(data["steps"]):
        card = add_round_box(slide, Inches(x + i * 1.8), Inches(2.95), Inches(1.55), Inches(2.3), colors[i % len(colors)])
        tf = card.text_frame
        tf.clear()
        tf.word_wrap = True
        p1 = tf.paragraphs[0]
        r1 = p1.add_run()
        r1.text = title
        r1.font.name = "Aptos Display"
        r1.font.size = Pt(16)
        r1.font.bold = True
        r1.font.color.rgb = TEXT
        p2 = tf.add_paragraph()
        p2.text = body
        p2.font.name = "Aptos"
        p2.font.size = Pt(11)
        p2.font.color.rgb = TEXT
        if i < len(data["steps"]) - 1:
            arrow = slide.shapes.add_textbox(Inches(x + i * 1.8 + 1.45), Inches(3.75), Inches(0.5), Inches(0.4))
            tf_arrow = arrow.text_frame
            p = tf_arrow.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = "→"
            run.font.size = Pt(24)
            run.font.bold = True
            run.font.color.rgb = SUBTEXT


def render_architecture(slide, data):
    add_decorations(slide)
    add_header(slide, data["title"], data["subtitle"], data.get("tag"))
    colors = [ROSE, MINT, BLUE, YELLOW]
    for i, (title, body) in enumerate(data["layers"]):
        y = 2.3 + i * 1.05
        card = add_round_box(slide, Inches(1.35), Inches(y), Inches(5.85), Inches(0.78), colors[i])
        tf = card.text_frame
        tf.clear()
        p1 = tf.paragraphs[0]
        r1 = p1.add_run()
        r1.text = title
        r1.font.name = "Aptos Display"
        r1.font.size = Pt(17)
        r1.font.bold = True
        r1.font.color.rgb = TEXT
        p2 = tf.add_paragraph()
        p2.text = body
        p2.font.name = "Aptos"
        p2.font.size = Pt(12)
        p2.font.color.rgb = SUBTEXT
        if i < len(data["layers"]) - 1:
            arr = slide.shapes.add_textbox(Inches(3.9), Inches(y + 0.72), Inches(0.8), Inches(0.35))
            tf_a = arr.text_frame
            p = tf_a.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = "↓"
            run.font.size = Pt(20)
            run.font.bold = True
            run.font.color.rgb = SUBTEXT
    add_phone(slide, 7.8, 2.55, 1.2, 2.4)


def render_split(slide, data):
    add_decorations(slide)
    add_header(slide, data["title"], data["subtitle"], data.get("tag"))
    left = add_round_box(slide, Inches(0.85), Inches(2.45), Inches(3.65), Inches(3.45), WHITE, LAVENDER)
    right = add_round_box(slide, Inches(4.85), Inches(2.45), Inches(3.65), Inches(3.45), WHITE, LAVENDER)
    add_card_text(left, data["left_title"], data["left_points"], MINT)
    add_card_text(right, data["right_title"], data["right_points"], PEACH)


def add_card_text(shape, title, points, pill_color):
    pill = shape.part.slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        shape.left + Inches(0.18),
        shape.top + Inches(0.2),
        Inches(1.55),
        Inches(0.36),
    )
    pill.fill.solid()
    pill.fill.fore_color.rgb = pill_color
    pill.line.color.rgb = pill_color
    add_text(pill, title, size=13, bold=True, align=PP_ALIGN.CENTER)
    pill.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    box = shape.text_frame
    box.clear()
    box.word_wrap = True
    box.margin_left = Inches(0.22)
    box.margin_right = Inches(0.15)
    box.margin_top = Inches(0.75)
    for i, point in enumerate(points):
        p = box.paragraphs[0] if i == 0 else box.add_paragraph()
        p.text = f"• {point}"
        p.font.name = "Aptos"
        p.font.size = Pt(18)
        p.font.color.rgb = TEXT
        p.space_after = Pt(8)


def render_domain(slide, data):
    add_decorations(slide)
    add_header(slide, data["title"], data["subtitle"], data.get("tag"))
    xs = [0.95, 3.35, 5.75]
    colors = [MINT, BLUE, PEACH]
    for i, (title, body) in enumerate(data["cards"]):
        add_card(slide, xs[i], 2.6, 2.0, 2.35, title, body, colors[i])
    mid = slide.shapes.add_textbox(Inches(4.45), Inches(5.2), Inches(1.2), Inches(0.35))
    tf = mid.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "User  ↔  Purchase  ↔  App"
    r.font.size = Pt(15)
    r.font.bold = True
    r.font.color.rgb = SUBTEXT
    footer = add_round_box(slide, Inches(1.25), Inches(5.8), Inches(6.9), Inches(0.55), YELLOW)
    add_text(footer, data["footer"], size=14, bold=True, align=PP_ALIGN.CENTER)


def render_formula(slide, data):
    add_decorations(slide)
    add_header(slide, data["title"], data["subtitle"], data.get("tag"))
    formula = add_round_box(slide, Inches(0.95), Inches(2.55), Inches(7.3), Inches(1.1), BLUE)
    add_text(formula, data["formula"], size=22, bold=True, align=PP_ALIGN.CENTER, font_name="Aptos Display")
    formula.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    points = add_round_box(slide, Inches(1.15), Inches(4.05), Inches(4.9), Inches(1.75), WHITE, LAVENDER)
    add_bullet_list(points, data["points"], size=16)
    add_avatar(slide, 7.2, 3.5, 1.1)


def render_flow(slide, data):
    add_decorations(slide)
    add_header(slide, data["title"], data["subtitle"], data.get("tag"))
    colors = [ROSE, MINT, BLUE, YELLOW, PEACH, LAVENDER]
    for i, step in enumerate(data["steps"]):
        y = 2.3 + i * 0.66
        box = add_round_box(slide, Inches(1.0 + (i % 2) * 0.35), Inches(y), Inches(6.95), Inches(0.5), colors[i % len(colors)])
        add_text(box, step, size=15, bold=(i == 0))
    add_phone(slide, 8.2, 2.85, 1.25, 2.55)


def render_api(slide, data):
    add_decorations(slide)
    add_header(slide, data["title"], data["subtitle"], data.get("tag"))
    colors = [MINT, BLUE, PEACH]
    for i, (api, purpose) in enumerate(data["apis"]):
        box = add_round_box(slide, Inches(0.95), Inches(2.45 + i * 1.15), Inches(7.35), Inches(0.82), colors[i % len(colors)])
        tf = box.text_frame
        tf.clear()
        p1 = tf.paragraphs[0]
        r1 = p1.add_run()
        r1.text = api
        r1.font.name = "Aptos Display"
        r1.font.size = Pt(16)
        r1.font.bold = True
        r1.font.color.rgb = TEXT
        p2 = tf.add_paragraph()
        p2.text = purpose
        p2.font.name = "Aptos"
        p2.font.size = Pt(12)
        p2.font.color.rgb = SUBTEXT
    add_phone(slide, 8.45, 2.65, 1.0, 2.0)


def render_testing(slide, data):
    add_decorations(slide)
    add_header(slide, data["title"], data["subtitle"], data.get("tag"))
    left = add_round_box(slide, Inches(0.9), Inches(2.55), Inches(3.5), Inches(3.2), WHITE, LAVENDER)
    right = add_round_box(slide, Inches(4.8), Inches(2.55), Inches(3.5), Inches(3.2), WHITE, LAVENDER)
    add_card_text(left, "Happy Paths", data["good"], MINT)
    add_card_text(right, "Edge Cases", data["bad"], PEACH)


def render_future(slide, data):
    add_decorations(slide)
    add_header(slide, data["title"], data["subtitle"], data.get("tag"))
    left = add_round_box(slide, Inches(0.95), Inches(2.5), Inches(3.45), Inches(3.25), WHITE, LAVENDER)
    right = add_round_box(slide, Inches(4.8), Inches(2.5), Inches(3.45), Inches(3.25), WHITE, LAVENDER)
    add_card_text(left, "Product Ideas", data["product"], YELLOW)
    add_card_text(right, "Technical Ideas", data["tech"], BLUE)


def render_closing(slide, data):
    add_decorations(slide)
    add_header(slide, data["title"], data["subtitle"], data.get("tag"))
    points = add_round_box(slide, Inches(0.95), Inches(2.5), Inches(6.0), Inches(2.6), WHITE, LAVENDER)
    add_bullet_list(points, data["points"], size=17)
    add_phone(slide, 7.5, 2.6, 1.45, 2.8)
    add_avatar(slide, 8.05, 5.05, 0.8)
    sticker = add_round_box(slide, Inches(6.95), Inches(5.8), Inches(1.8), Inches(0.5), YELLOW)
    add_text(sticker, "Thank you!", size=16, bold=True, align=PP_ALIGN.CENTER)


def build():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    renderers = {
        "cover": render_cover,
        "hero": render_hero,
        "journey": render_journey,
        "architecture": render_architecture,
        "split": render_split,
        "domain": render_domain,
        "formula": render_formula,
        "flow": render_flow,
        "api": render_api,
        "testing": render_testing,
        "future": render_future,
        "closing": render_closing,
    }

    for slide_data in SLIDES:
        slide = prs.slides.add_slide(blank)
        set_bg(slide)
        renderers[slide_data["kind"]](slide, slide_data)

    prs.save(OUTPUT)


if __name__ == "__main__":
    build()
